"""Generated document fixtures for parser tests."""

from __future__ import annotations

from io import BytesIO
from typing import TYPE_CHECKING, cast

from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from blend_brain.document_ingestion.domain import DocumentSource

if TYPE_CHECKING:
    from pathlib import Path

    from pypdf.generic import EncodedStreamObject


def source(filename: str, content: bytes) -> DocumentSource:
    """Build a deterministic in-memory source for adapter tests."""
    return DocumentSource(
        source_id=f"memory://{filename}",
        filename=filename,
        content=content,
        size_bytes=len(content),
        sha256="a" * 64,
    )


def make_docx() -> bytes:
    """Create a representative Word document."""
    document = Document()
    document.core_properties.title = "Word Project"
    document.core_properties.author = "Blend"
    document.add_paragraph("Executive summary", style="Heading 1")
    document.add_paragraph("A reusable analytics platform.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Client"
    table.cell(0, 1).text = "Industry"
    table.cell(1, 0).text = "Example Co"
    table.cell(1, 1).text = "Retail"
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def make_pptx() -> bytes:
    """Create a representative presentation with text, a table, and notes."""
    presentation = Presentation()
    presentation.core_properties.title = "Project Deck"
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    if title is not None:
        title.text = "Architecture"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(1))
    text_box.text = "Composable ingestion"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Layer"
    table.cell(0, 1).text = "Owner"
    table.cell(1, 0).text = "Domain"
    table.cell(1, 1).text = "Platform"
    notes_frame = slide.notes_slide.notes_text_frame
    if notes_frame is not None:
        notes_frame.text = "Discuss security limits"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def make_pdf(*, include_blank_page: bool = False, encrypted: bool = False) -> bytes:
    """Create a digital PDF whose text is extractable by pypdf."""
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    page[NameObject("/Resources")] = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject(
                {
                    NameObject("/F1"): DictionaryObject(
                        {
                            NameObject("/Type"): NameObject("/Font"),
                            NameObject("/Subtype"): NameObject("/Type1"),
                            NameObject("/BaseFont"): NameObject("/Helvetica"),
                        }
                    )
                }
            )
        }
    )
    stream = DecodedStreamObject()
    stream.set_data(b"BT /F1 12 Tf 72 720 Td (Project knowledge) Tj ET")
    page.replace_contents(cast("EncodedStreamObject", stream))
    if include_blank_page:
        writer.add_blank_page(width=612, height=792)
    writer.add_metadata({"/Title": "PDF Project", "/Author": "Blend"})
    if encrypted:
        writer.encrypt("password")
    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def write(path: Path, content: bytes) -> Path:
    """Write generated fixture bytes using pytest-owned temporary storage."""
    path.write_bytes(content)
    return path
