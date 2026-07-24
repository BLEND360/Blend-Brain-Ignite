"""PowerPoint Open XML parser adapter."""

from __future__ import annotations

from io import BytesIO
from typing import Any
from zipfile import BadZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from blend_brain.document_ingestion.domain import (
    CorruptDocumentError,
    DocumentFormat,
    DocumentMetadata,
    DocumentSection,
    DocumentSource,
    ExtractedDocument,
    SectionKind,
    SectionLocator,
)
from blend_brain.document_ingestion.infrastructure.parsers.common import (
    build_document,
    normalize_datetime,
    optional_text,
)


class PptxParser:
    """Extract ordered slide text, tables, and speaker notes."""

    @property
    def document_format(self) -> DocumentFormat:
        """Return the supported format."""
        return DocumentFormat.PPTX

    def parse(self, source: DocumentSource) -> ExtractedDocument:
        """Create one citation section per presentation slide."""
        try:
            presentation = Presentation(BytesIO(source.content))
            sections: list[DocumentSection] = []
            for slide_number, slide in enumerate(presentation.slides, start=1):
                lines: list[str] = []
                shapes = sorted(slide.shapes, key=lambda shape: (shape.top, shape.left))
                for shape in shapes:
                    lines.extend(self._shape_lines(shape))
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    notes = notes_frame.text.strip() if notes_frame is not None else ""
                    if notes:
                        lines.extend(("Speaker notes", notes))
                sections.append(
                    DocumentSection(
                        sequence=slide_number,
                        kind=SectionKind.SLIDE,
                        text="\n".join(lines),
                        locator=SectionLocator(slide_number=slide_number),
                    )
                )
            properties = presentation.core_properties
            metadata = DocumentMetadata(
                title=optional_text(properties.title),
                author=optional_text(properties.author),
                subject=optional_text(properties.subject),
                created_at=normalize_datetime(properties.created),
                modified_at=normalize_datetime(properties.modified),
            )
            return build_document(source, self.document_format, sections, metadata=metadata)
        except (BadZipFile, ValueError, KeyError, TypeError) as exception:
            raise CorruptDocumentError(
                "PowerPoint document could not be parsed",
                filename=source.filename,
            ) from exception

    def _shape_lines(self, shape: Any) -> list[str]:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            lines: list[str] = []
            for child in sorted(shape.shapes, key=lambda item: (item.top, item.left)):
                lines.extend(self._shape_lines(child))
            return lines
        if getattr(shape, "has_table", False):
            return [
                " | ".join(cell.text.strip() for cell in row.cells)
                for row in shape.table.rows
                if any(cell.text.strip() for cell in row.cells)
            ]
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            return [text] if text else []
        return []
